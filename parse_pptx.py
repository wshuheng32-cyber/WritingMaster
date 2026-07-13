#!/usr/bin/env python3
"""Parse the grasping PPT part 3 to understand what the teacher wants."""

from pptx import Presentation
import sys

pptx_path = '/home/ros/ros2_ws/机械臂自主抓取-第三部分.pptx'
prs = Presentation(pptx_path)

print(f'幻灯片数量: {len(prs.slides)}')
print('=' * 80)

for i, slide in enumerate(prs.slides):
    print(f'\n===== 第 {i+1} 页 =====')
    print(f'布局: {slide.slide_layout.name}')

    # Extract all text
    for shape in slide.shapes:
        if hasattr(shape, 'text') and shape.text.strip():
            print(f'  [{shape.shape_type}] {shape.text[:800]}')
        if shape.has_table:
            table = shape.table
            print('  --- 表格内容 ---')
            for row in table.rows:
                cells = [cell.text for cell in row.cells]
                print(f'    | {" | ".join(cells)}')
    print()